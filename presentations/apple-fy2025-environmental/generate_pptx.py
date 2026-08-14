#!/usr/bin/env python3
"""Generate Apple FY2025 Environmental Progress PowerPoint deck."""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

OUTPUT = Path(__file__).parent / "Apple_FY2025_Environmental_Progress.pptx"

APPLE_GREEN = RGBColor(0x1A, 0x7F, 0x37)
APPLE_GREEN_BRIGHT = RGBColor(0x30, 0xD1, 0x58)
DARK = RGBColor(0x1D, 0x1D, 0x1F)
GRAY = RGBColor(0x6E, 0x6E, 0x73)
LIGHT_BG = RGBColor(0xF5, 0xF5, 0xF7)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def set_slide_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, size=18, bold=False, color=DARK, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return box


def add_footer(slide, text):
    add_textbox(slide, Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.35),
                 text, size=9, color=GRAY, align=PP_ALIGN.CENTER)


def add_speaker_notes(slide, notes):
    slide.notes_slide.notes_text_frame.text = notes


def slide_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, LIGHT_BG)

    add_textbox(slide, Inches(1), Inches(1.2), Inches(11.3), Inches(0.4),
                "APPLE 2030", size=12, bold=True, color=APPLE_GREEN, align=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(0.8), Inches(2.0), Inches(11.7), Inches(2.2),
                "Apple FY2025:\nEnvironmental Progress as a Business & Sales Driver",
                size=36, bold=True, color=DARK, align=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(1), Inches(4.5), Inches(11.3), Inches(0.6),
                "Decoupling Business Growth from Emissions",
                size=22, color=GRAY, align=PP_ALIGN.CENTER)

    add_footer(slide, "Confidential | Board of Directors & Sales Leadership Strategy")

    add_speaker_notes(slide,
        "Open with the strategic framing — environmental leadership is not a cost center "
        "but a revenue and risk-mitigation engine for premium positioning and enterprise procurement.")


def slide_snapshot(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_textbox(slide, Inches(0.7), Inches(0.5), Inches(12), Inches(0.8),
                "FY25 Snapshot: Turning Sustainability into Shareholder Value",
                size=28, bold=True, color=DARK)

    stats = [
        ("60%", "Reduction in gross emissions vs. 2015 baseline — sustained through business growth."),
        ("30%", "Of total materials shipped are now recycled or renewable."),
        ("100%", "Plastic-free product packaging — fully fiber-based."),
        ("$4.7B", "Issued in green bonds to fund supply chain transformation."),
    ]

    positions = [(0.7, 1.6), (6.8, 1.6), (0.7, 4.0), (6.8, 4.0)]
    for (left, top), (stat, desc) in zip(positions, stats):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top),
                                       Inches(5.8), Inches(2.1))
        shape.fill.solid()
        shape.fill.fore_color.rgb = LIGHT_BG
        shape.line.fill.background()

        add_textbox(slide, Inches(left + 0.3), Inches(top + 0.25), Inches(5.2), Inches(0.9),
                    stat, size=44, bold=True, color=APPLE_GREEN)
        add_textbox(slide, Inches(left + 0.3), Inches(top + 1.1), Inches(5.2), Inches(0.8),
                    desc, size=14, color=DARK)

    add_speaker_notes(slide,
        "This year marks a critical inflection point. We have proven that aggressive environmental "
        "targets do not stifle growth. In fact, a 60% reduction in gross emissions has occurred "
        "alongside sustained revenue growth. These metrics are now key differentiators in both "
        "consumer retail and enterprise procurement.")


def slide_derisk(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_textbox(slide, Inches(0.7), Inches(0.5), Inches(6), Inches(1.2),
                "De-risking the Business:\nProgress Toward Total Carbon Neutrality",
                size=24, bold=True, color=DARK)

    bullets = [
        "Corporate Operations: 100% Carbon Neutral (Achieved)",
        "Full Value Chain: 60% progress toward 2030 neutrality",
        "Data Centers: 100% certified to AWS Standard (8/8 facilities)",
    ]
    top = 2.0
    for b in bullets:
        add_textbox(slide, Inches(0.9), Inches(top), Inches(5.5), Inches(0.6), f"•  {b}", size=15, color=DARK)
        top += 0.7

    # Progress rings (simplified as circles with labels)
    for pct, label, x in [(100, "Operations", 8.5), (60, "Value Chain", 10.5)]:
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(2.5), Inches(1.6), Inches(1.6))
        circle.fill.background()
        circle.line.color.rgb = APPLE_GREEN if pct == 100 else APPLE_GREEN_BRIGHT
        circle.line.width = Pt(8)
        add_textbox(slide, Inches(x), Inches(3.0), Inches(1.6), Inches(0.5),
                    f"{pct}%", size=22, bold=True, color=DARK, align=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(x), Inches(4.2), Inches(1.6), Inches(0.4),
                    label, size=11, color=GRAY, align=PP_ALIGN.CENTER)

    add_speaker_notes(slide,
        "We have eliminated direct operational carbon risk. Our focus is now on the value chain. "
        "Achieving 60% progress toward 2030 neutrality ensures we are ahead of impending global "
        "regulatory frameworks regarding Scope 3 emissions reporting.")


def slide_premium(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_textbox(slide, Inches(0.7), Inches(0.5), Inches(6), Inches(1.2),
                'Premium Products, Premium Ethics:\nThe "Recycled" Advantage',
                size=24, bold=True, color=DARK)

    bullets = [
        "100% Recycled: Cobalt, Gold, Tin, and Rare Earth Elements in select components and batteries.",
        "The Unboxing Experience: 100% fiber-based packaging — plastics completely eliminated.",
    ]
    top = 2.2
    for b in bullets:
        add_textbox(slide, Inches(0.9), Inches(top), Inches(5.8), Inches(0.8), f"•  {b}", size=14, color=DARK)
        top += 0.9

    # Device outline
    device = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.0), Inches(1.8), Inches(2.5), Inches(4.0))
    device.fill.solid()
    device.fill.fore_color.rgb = RGBColor(0xD2, 0xD2, 0xD7)
    device.line.fill.background()

    for label, y in [("Battery\n100% recycled cobalt", 2.2), ("Logic Board\n100% recycled tin & gold", 3.8),
                     ("Magnets\n100% recycled rare earths", 1.5)]:
        add_textbox(slide, Inches(10.8), Inches(y), Inches(2.2), Inches(0.7), label, size=10, color=DARK)

    add_speaker_notes(slide,
        "Enterprise clients are under immense pressure to reduce their own Scope 3 emissions. "
        "By purchasing Apple hardware made with 100% recycled key minerals, our B2B clients "
        "instantly lower their procurement footprint. This is a powerful closing tool for our "
        "enterprise sales teams.")


def slide_efficiency(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_textbox(slide, Inches(0.7), Inches(0.5), Inches(12), Inches(0.8),
                "Operational Efficiency: Stewarding Resources to Protect Margins",
                size=26, bold=True, color=DARK)

    cards = [
        ("17B", "Gallons of freshwater saved across the supply chain in FY25."),
        ("100%", "Zero waste to landfill at all final assembly sites."),
        ("55%", "Progress toward our 2030 freshwater replenishment goal."),
    ]

    for i, (stat, desc) in enumerate(cards):
        left = 0.7 + i * 4.2
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(2.0),
                                       Inches(3.8), Inches(3.5))
        shape.fill.solid()
        shape.fill.fore_color.rgb = LIGHT_BG
        shape.line.fill.background()

        add_textbox(slide, Inches(left + 0.3), Inches(2.8), Inches(3.2), Inches(0.8),
                    stat, size=40, bold=True, color=APPLE_GREEN, align=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(left + 0.3), Inches(3.8), Inches(3.2), Inches(1.2),
                    desc, size=13, color=DARK, align=PP_ALIGN.CENTER)

    add_speaker_notes(slide,
        "Water scarcity and waste disposal are significant operational risks and costs in "
        "manufacturing regions. Saving 17 billion gallons and achieving zero waste at final "
        "assembly sites secures our supply chain against climate disruptions and reduces utility overhead.")


def slide_emissions(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_textbox(slide, Inches(0.7), Inches(0.5), Inches(6), Inches(1.0),
                "Targeting Emissions Where They Matter to the Customer",
                size=24, bold=True, color=DARK)

    items = [
        ("53%", "Manufacturing (Scope 3)", True),
        ("27%", "Product Use (Scope 3 — Customer devices)", False),
        ("16%", "Product Transport (Scope 3)", False),
        ("4%", "Business Operations (Scopes 1–2)", False),
    ]
    top = 1.8
    for pct, label, highlight in items:
        size = 16 if highlight else 14
        color = DARK
        add_textbox(slide, Inches(0.9), Inches(top), Inches(5.5), Inches(0.5),
                    f"●  {pct}  {label}", size=size, bold=highlight, color=color)
        top += 0.65

    # Donut placeholder
    outer = slide.shapes.add_shape(MSO_SHAPE.DONUT, Inches(8.0), Inches(1.8), Inches(3.5), Inches(3.5))
    outer.fill.solid()
    outer.fill.fore_color.rgb = APPLE_GREEN
    outer.line.fill.background()

    add_textbox(slide, Inches(8.5), Inches(3.0), Inches(2.5), Inches(0.8),
                "Scope 3\n96%", size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_speaker_notes(slide,
        '27% of our footprint comes from "Product Use"—the energy our customers consume. '
        "Through Apple Silicon efficiency, we are drastically lowering the energy bills and "
        "downstream emissions for our users. We don't just sell a device; we sell long-term energy efficiency.")


def slide_funding(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_textbox(slide, Inches(0.7), Inches(0.5), Inches(6), Inches(1.0),
                "Funding the Transition: ROI on Environmental Programs",
                size=24, bold=True, color=DARK)

    add_textbox(slide, Inches(0.7), Inches(1.8), Inches(5.5), Inches(1.5),
                "43M", size=72, bold=True, color=APPLE_GREEN)
    add_textbox(slide, Inches(0.7), Inches(3.2), Inches(5.5), Inches(0.4),
                "Metric Tons", size=18, bold=True, color=DARK)
    add_textbox(slide, Inches(0.7), Inches(3.7), Inches(5.5), Inches(0.5),
                "Estimated emissions avoided by FY25 programs.", size=14, color=GRAY)

    add_textbox(slide, Inches(0.7), Inches(4.6), Inches(5.5), Inches(0.7),
                "$4.7B", size=36, bold=True, color=APPLE_GREEN)
    add_textbox(slide, Inches(0.7), Inches(5.3), Inches(5.5), Inches(0.4),
                "Successfully issued in green bonds.", size=14, color=GRAY)

    cert = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.5), Inches(2.0), Inches(4.5), Inches(3.5))
    cert.fill.solid()
    cert.fill.fore_color.rgb = RGBColor(0xF0, 0xFA, 0xF2)
    cert.line.color.rgb = APPLE_GREEN
    cert.line.width = Pt(2)
    add_textbox(slide, Inches(8.5), Inches(4.5), Inches(2.5), Inches(0.4),
                "GREEN BOND", size=12, bold=True, color=APPLE_GREEN, align=PP_ALIGN.CENTER)

    add_speaker_notes(slide,
        "Our environmental strategy is heavily self-funded and leveraged through favorable green "
        "finance instruments. The $4.7B in green bonds not only funds necessary infrastructure "
        "but also strengthens our ESG ratings, appealing to institutional investors and lowering our cost of capital.")


def slide_action(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_textbox(slide, Inches(0.7), Inches(0.5), Inches(12), Inches(0.8),
                "Leveraging FY25 Wins in the Market (Action Plan)",
                size=26, bold=True, color=DARK)

    steps = [
        ("1", "Enterprise Pitch Decks",
         'Integrate "100% Recycled Cobalt/Gold" and "Scope 3 Reduction" stats into B2B sales materials.'),
        ("2", "Retail Training",
         'Equip Apple Store teams with "17 Billion Gallons Saved" talking points for customer engagement.'),
        ("3", "Marketing",
         "Launch a campaign highlighting the 100% plastic-free unboxing experience as a premium feature."),
    ]

    for i, (num, title, desc) in enumerate(steps):
        left = 0.7 + i * 4.2
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(2.0),
                                       Inches(3.8), Inches(3.8))
        shape.fill.solid()
        shape.fill.fore_color.rgb = LIGHT_BG
        shape.line.fill.background()

        num_shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(left + 0.3), Inches(2.3),
                                           Inches(0.5), Inches(0.5))
        num_shape.fill.solid()
        num_shape.fill.fore_color.rgb = APPLE_GREEN
        num_shape.line.fill.background()
        add_textbox(slide, Inches(left + 0.3), Inches(2.35), Inches(0.5), Inches(0.4),
                    num, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        add_textbox(slide, Inches(left + 0.3), Inches(3.0), Inches(3.2), Inches(0.5),
                    title, size=16, bold=True, color=DARK)
        add_textbox(slide, Inches(left + 0.3), Inches(3.6), Inches(3.2), Inches(1.5),
                    desc, size=12, color=GRAY)

    add_speaker_notes(slide,
        "We are directing Sales and Marketing to immediately integrate these FY25 wins into their "
        "Q3 and Q4 go-to-market strategies. We need to empower our enterprise reps to use our "
        "recycled material stats to win RFPs, and equip our retail teams to connect emotionally "
        "with consumers. Thank you.")


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_title(prs)
    slide_snapshot(prs)
    slide_derisk(prs)
    slide_premium(prs)
    slide_efficiency(prs)
    slide_emissions(prs)
    slide_funding(prs)
    slide_action(prs)

    prs.save(str(OUTPUT))
    print(f"Generated: {OUTPUT}")


if __name__ == "__main__":
    main()
